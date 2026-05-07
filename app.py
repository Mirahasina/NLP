import streamlit as st
import spacy
import nltk
from nltk.corpus import stopwords
from collections import Counter
import heapq
import re
import os
import pickle
import pandas as pd
import PyPDF2
from docx import Document
from pptx import Presentation
import io
import numpy as np
from sklearn.metrics.pairwise import cosine_similarity
try:
    from sentence_transformers import SentenceTransformer
except ImportError:
    SentenceTransformer = None

import subprocess
import sys

def install_model(model):
    subprocess.run([sys.executable, "-m", "spacy", "download", model],
                   stdout=subprocess.DEVNULL,
                   stderr=subprocess.DEVNULL)

@st.cache_resource
def load_models():
    try:
        nlp_en = spacy.load("en_core_web_sm")
    except OSError:
        install_model("en_core_web_sm")
        nlp_en = spacy.load("en_core_web_sm")

    try:
        nlp_fr = spacy.load("fr_core_news_sm")
    except OSError:
        install_model("fr_core_news_sm")
        nlp_fr = spacy.load("fr_core_news_sm")

    return nlp_en, nlp_fr

nlp_en, nlp_fr = load_models()

st.title("NLP FR / EN App")

st.set_page_config(
    page_title="Résumeur IA Multilingue",
    layout="centered"
)

st.markdown("""
<style>
    .stApp {
        background: linear-gradient(135deg, #0f172a 0%, #1e293b 100%);
        color: #f8fafc;
    }
    .main-title {
        font-size: 3rem !important;
        font-weight: 700 !important;
        background: linear-gradient(to right, #38bdf8, #818cf8);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        text-align: center;
        margin-bottom: 0.5rem;
    }
    .sub-title {
        text-align: center;
        color: #94a3b8;
        font-size: 1.1rem;
        margin-bottom: 2rem;
    }
    .stTextArea textarea {
        background-color: #1e293b !important;
        color: #f8fafc !important;
        border: 1px solid #334155 !important;
        border-radius: 12px !important;
    }
    .stButton>button {
        width: 100%;
        background: linear-gradient(90deg, #38bdf8 0%, #818cf8 100%) !important;
        color: white !important;
        border: none !important;
        padding: 0.75rem !important;
        font-weight: 600 !important;
        border-radius: 12px !important;
        transition: all 0.3s ease !important;
        box-shadow: 0 4px 6px -1px rgba(0, 0, 0, 0.1), 0 2px 4px -1px rgba(0, 0, 0, 0.06);
    }
    .stButton>button:hover {
        transform: translateY(-2px);
        box-shadow: 0 10px 15px -3px rgba(0, 0, 0, 0.1), 0 4px 6px -2px rgba(0, 0, 0, 0.05);
        opacity: 0.9;
    }
    .summary-box {
        background-color: #1e293b;
        padding: 2rem;
        border-radius: 16px;
        border-left: 5px solid #38bdf8;
        margin-top: 2rem;
    }
    .metric-card {
        background: rgba(255, 255, 255, 0.05);
        padding: 1rem;
        border-radius: 8px;
        text-align: center;
    }
</style>
""", unsafe_allow_html=True)

def extract_text_from_pdf(file):
    pdf_reader = PyPDF2.PdfReader(file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text() + "\n"
    return text

def extract_text_from_docx(file):
    doc = Document(file)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

class SummarizerFinal:
    def __init__(self):
        self.idf_weights = {'en': {}, 'fr': {}}

    def summarize(self, text, n=3, nlp=None, stop_words=None, threshold=0.8):
        if not text or len(text) < 10: return ""
        if not nlp: return ""
        
        lang_code = 'en' if nlp.meta['lang'] == 'en' else 'fr'
        idfs = self.idf_weights.get(lang_code, {})
        
        doc = nlp(text)
        sentences = [sent for sent in doc.sents if len(sent.text.split()) >= 8]
        if not sentences: return ""
        
        scored_sentences = []
        for sent in sentences:
            text_s = sent.text.strip()
            if re.match(r'^(\d+|[IVX]+)[\.\)]', text_s) or text_s.endswith(':') or '%' in text_s or len(text_s) < 20:
                score = 0
            else:
                score = sum(idfs.get(t.lemma_.lower(), 1.0) for t in sent if not t.is_stop and t.pos_ in ['NOUN', 'ADJ', 'VERB'])
            
            score = score / (len(sent)**0.7 + 1)
            scored_sentences.append((sent, score))
        
        scored_sentences.sort(key=lambda x: x[1], reverse=True)
        
        final_sentences = []
        selected_texts = []
        
        for sent, score in scored_sentences:
            if len(final_sentences) >= n: break
            if score <= 0: continue
            
            is_redundant = False
            for prev_text in selected_texts:
                s1 = set(sent.text.lower().split())
                s2 = set(prev_text.lower().split())
                if not s1 or not s2: continue
                overlap = len(s1 & s2) / min(len(s1), len(s2))
                if overlap > threshold:
                    is_redundant = True
                    break
            
            if not is_redundant:
                final_sentences.append(sent)
                selected_texts.append(sent.text)
        
        res = sorted(final_sentences, key=lambda x: x.start)
        return "\n- " + "\n- ".join([s.text.strip() for s in res])

@st.cache_resource
def load_resources(lang_code):
    try:
        nltk.data.find('corpora/stopwords')
    except LookupError:
        nltk.download('stopwords')
    
    models_to_try = [
        'en_core_web_sm' if lang_code == 'en' else 'fr_core_news_sm',
        'en_core_web_sm'
    ]
    
    nlp = None
    for model_name in models_to_try:
        try:
            nlp = spacy.load(model_name, disable=['ner', 'textcat'])
            break
        except OSError:
            continue
            
    if nlp is None:
        os.system("python3 -m spacy download en_core_web_sm")
        nlp = spacy.load('en_core_web_sm', disable=['ner', 'textcat'])
        
    stop_words = set(stopwords.words('english' if lang_code == 'en' else 'french'))
    return nlp, stop_words

@st.cache_resource
def load_trained_model():
    if os.path.exists('modele_final.pkl'):
        try:
            with open('modele_final.pkl', 'rb') as f:
                return pickle.load(f)
        except Exception:
            return SummarizerFinal()
    return SummarizerFinal()

@st.cache_resource
def load_bert_model():
    if SentenceTransformer is None:
        return None
    return SentenceTransformer('paraphrase-multilingual-MiniLM-L12-v2')

def summarize_bert(text, n=3, model_bert=None, threshold=0.8):
    if not text or not model_bert: return ""
    
    sentences = re.split(r'(?<=[.!?]) +', text)
    sentences = [s.strip() for s in sentences if len(s.split()) > 5]
    if not sentences: return ""
    
    embeddings = model_bert.encode(sentences)
    doc_embedding = embeddings.mean(axis=0).reshape(1, -1)
    
    similarities = cosine_similarity(embeddings, doc_embedding).flatten()
    scored_sentences = sorted(zip(sentences, similarities, range(len(sentences))), key=lambda x: x[1], reverse=True)
    
    final_sentences = []
    selected_embeddings = []
    
    for sent, score, idx in scored_sentences:
        if len(final_sentences) >= n:
            break
            
        is_redundant = False
        if selected_embeddings and threshold < 1.0:
            current_emb = embeddings[idx].reshape(1, -1)
            sims = cosine_similarity(current_emb, np.array(selected_embeddings))
            if sims.max() > threshold:
                is_redundant = True
        
        if not is_redundant:
            final_sentences.append((sent, idx))
            selected_embeddings.append(embeddings[idx])
            
    return " ".join([s[0] for s in sorted(final_sentences, key=lambda x: x[1])])

model = load_trained_model()

def main():
    st.markdown('<h1 class="main-title">Resumeur IA</h1>', unsafe_allow_html=True)
    st.markdown('<p class="sub-title">Analyse semantique multilingue et Extraction d\'informations</p>', unsafe_allow_html=True)

    # --- SIDEBAR ---
    lang = st.sidebar.radio("Langue / Language", ["Francais", "English"])
    lang_code = 'fr' if lang == "Francais" else 'en'
    
    nlp_active, stop_words_active = load_resources(lang_code)

    st.sidebar.markdown("---")
    st.sidebar.subheader("Configuration")
    model_type = st.sidebar.selectbox("Moteur de resume", ["BERT (Semantique)", "Automatique (TF-IDF)"])
    
    with st.sidebar.expander("Parametres avances"):
        use_filter = st.checkbox("Filtre de redondance", value=True)
        redundancy_threshold = st.slider("Seuil de filtrage", 0.5, 1.0, 0.7) if use_filter else 1.0
        num_sentences = st.slider("Nombre de phrases", 1, 20, 3)
        
    if model_type == "BERT (Semantique)":
        with st.sidebar:
            with st.spinner("Chargement de BERT..."):
                bert_model = load_bert_model()
                if bert_model is None:
                    model_type = "Automatique (TF-IDF)"
    
    # --- MAIN UI ---
    tab1, tab2 = st.tabs(["Saisie et Resume", "Analyses Detaillees"])

    labels = {
        'fr': {
            'upload': "Importer un document (.pdf, .docx, .pptx, .txt, .csv)",
            'text_area': "Contenu du document :",
            'placeholder': "Le texte apparaitra ici...",
            'button': "Generer le resume",
            'spinner': "Analyse en cours...",
            'result': "Resultat du Resume",
            'stats': "Statistiques",
            'keywords': "Mots-cles principaux",
            'download': "Telecharger le resume (.txt)"
        },
        'en': {
            'upload': "Upload a document (.pdf, .docx, .pptx, .txt, .csv)",
            'text_area': "Document content:",
            'placeholder': "Text will appear here...",
            'button': "Generate Summary",
            'spinner': "Analyzing...",
            'result': "Summary Result",
            'stats': "Statistics",
            'keywords': "Main Keywords",
            'download': "Download summary (.txt)"
        }
    }[lang_code]

    with tab1:
        uploaded_file = st.file_uploader(labels['upload'], type=["pdf", "docx", "pptx", "txt", "csv"])
        input_text = ""
        
        if uploaded_file is not None:
            try:
                if uploaded_file.name.endswith('.pdf'):
                    input_text = extract_text_from_pdf(uploaded_file)
                elif uploaded_file.name.endswith('.docx'):
                    input_text = extract_text_from_docx(uploaded_file)
                elif uploaded_file.name.endswith('.pptx'):
                    input_text = extract_text_from_pptx(uploaded_file)
                elif uploaded_file.name.endswith('.txt'):
                    input_text = uploaded_file.read().decode("utf-8")
                elif uploaded_file.name.endswith('.csv'):
                    df = pd.read_csv(uploaded_file)
                    col_name = st.selectbox("Colonne a resumer", df.columns)
                    input_text = "\n".join(df[col_name].astype(str).tolist()[:50])
            except Exception as e:
                st.error(f"Erreur de lecture : {e}")

        input_text = st.text_area(labels['text_area'], value=input_text, height=250, placeholder=labels['placeholder'])
        
        if st.button(labels['button'], use_container_width=True):
            if not input_text.strip():
                st.warning("Veuillez fournir du texte.")
            else:
                with st.spinner(labels['spinner']):
                    clean_text = re.sub(' +', ' ', input_text.replace('\n', ' ')).strip()
                    try:
                        progress_bar = st.progress(0)
                        for i in range(100):
                            import time
                            time.sleep(0.005)
                            progress_bar.progress(i + 1)
                            
                        if model_type == "BERT (Semantique)":
                            summary = summarize_bert(clean_text, n=num_sentences, model_bert=bert_model, threshold=redundancy_threshold)
                        else:
                            summary = model.summarize(clean_text, n=num_sentences, nlp=nlp_active, stop_words=stop_words_active, threshold=redundancy_threshold)
                        
                        st.session_state['last_summary'] = summary
                        st.session_state['last_input'] = input_text
                    except Exception as e:
                        st.error(f"Erreur : {e}")
                        summary = ""

                if summary:
                    st.markdown('<div class="summary-box">', unsafe_allow_html=True)
                    st.subheader(labels['result'])
                    st.write(summary)
                    st.markdown('</div>', unsafe_allow_html=True)
                    st.download_button(labels['download'], summary, file_name="resume.txt")

    with tab2:
        if 'last_summary' in st.session_state and st.session_state['last_summary']:
            s_text = st.session_state['last_summary']
            i_text = st.session_state['last_input']
            
            st.subheader("Comparaison")
            c1, c2 = st.columns(2)
            with c1:
                st.info("Original")
                st.write(i_text[:500] + "...")
            with c2:
                st.success("Resume")
                st.write(s_text)
            
            st.divider()
            
            st.subheader(labels['stats'])
            m1, m2, m3 = st.columns(3)
            w_in = len(i_text.split())
            w_out = len(s_text.split())
            reduction = 100 - (w_out/w_in*100) if w_in > 0 else 0
            
            m1.metric("Mots Entree", w_in)
            m2.metric("Mots Sortie", w_out)
            m3.metric("Reduction", f"{reduction:.1f}%")
            
            st.divider()
            
            st.subheader(labels['keywords'])
            doc = nlp_active(i_text.lower())
            words = [token.text for token in doc if not token.is_stop and not token.is_punct and token.pos_ in ['NOUN', 'PROPN']]
            common = Counter(words).most_common(10)
            
            cols = st.columns(5)
            for i, (word, count) in enumerate(common):
                cols[i%5].markdown(f"**{word}** ({count})")
        else:
            st.info("Generez un resume pour voir les analyses.")

if __name__ == "__main__":
    main()
